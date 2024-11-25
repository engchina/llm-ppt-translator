import os
import time

import gradio as gr
import oci
from dotenv import load_dotenv, find_dotenv
from openai import OpenAI
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE

# read local .env file
_ = load_dotenv(find_dotenv())

# Ensure the 'outputs' folder exists
output_dir = os.path.join(os.path.dirname(__file__), "outputs")
os.makedirs(output_dir, exist_ok=True)

client = OpenAI(api_key=os.environ["OPENAI_API_KEY"], base_url=os.environ["OPENAI_BASE_URL"])


def translate_text(text, target_lang, model_name="gpt-4"):
    max_attempts = 5  # 最大尝试次数

    if model_name == "gpt-4":
        for attempt in range(max_attempts):
            try:
                completion = client.chat.completions.create(
                    model=os.environ["OPENAI_MODEL_NAME"],
                    # messages=[
                    #     {"role": "system", "content": "You are a helpful assistant that translates text."},
                    #     {"role": "user",
                    #      "content": f"Translate to {target_lang}, maintain the original tone and style, DO NOT translate or modify placeholders in the format [PLACEHOLDER_X]. Ensure the placeholders remain in their original positions and with the same quantity as in the original text. Only output the translated text. \n\nText: {text}"}
                    # ]
                    messages=[
                        {"role": "system",
                         "content": "You are a translation assistant specialized in maintaining the concise and professional style often used in presentation slides. Your task is to translate text while preserving placeholders and respecting language-specific style conventions."},
                        {"role": "user",
                         "content": f"""
                             Translate the text below into {target_lang}, adhering to the following rules:
                             1. Keep the original tone and style, ensuring the translation is concise and suitable for presentation slides.
                             2. Avoid overly formal or verbose expressions. For example:
                                - In Japanese, avoid 'です' and 'ます' unless absolutely necessary.
                                - In Chinese, use straightforward and professional wording.
                                - In English, prioritize brevity and clarity.
                             3. Do not translate or modify placeholders in the format [PLACEHOLDER_X]. Ensure placeholders remain in their original positions and with the same quantity as in the original text.
                             4. Only output the translated text without explanations or extra comments.
            
                             Text: {text}
                         """}
                    ]
                )
                print(f"{text=}")
                print(f"translated: {completion.choices[0].message.content}\n")
                return completion.choices[0].message.content
            except Exception as e:
                if attempt < max_attempts - 1:
                    print(f"Attempt {attempt + 1} failed with error: {str(e)}. Retrying...")
                    time.sleep(1)
                else:
                    print(f"All {max_attempts} attempts failed with error: {str(e)}")
                    return text
    else:
        for attempt in range(max_attempts):
            try:
                # Initialize OCI GenAI client
                config = oci.config.from_file('~/.oci/config', os.environ.get("CONFIG_PROFILE"))
                endpoint = "https://inference.generativeai.us-chicago-1.oci.oraclecloud.com"

                generative_ai_inference_client = oci.generative_ai_inference.GenerativeAiInferenceClient(
                    config=config,
                    service_endpoint=endpoint,
                    retry_strategy=oci.retry.NoneRetryStrategy(),
                    timeout=(10, 240)
                )

                # Prepare chat request
                chat_detail = oci.generative_ai_inference.models.ChatDetails()
                chat_request = oci.generative_ai_inference.models.CohereChatRequest()

                # chat_request.preamble_override = "You are a helpful assistant that translates text."
                # chat_request.message = f"Translate to {target_lang}, maintain the original tone and style, DO NOT translate or modify placeholders in the format [PLACEHOLDER_X]. Ensure the placeholders remain in their original positions and with the same quantity as in the original text. Only output the translated text. \n\nText: {text}"
                chat_request.preamble_override = "You are a translation assistant specialized in maintaining the concise and professional style often used in presentation slides. Your task is to translate text while preserving placeholders and respecting language-specific style conventions."
                chat_request.message = f"""
                    Translate the text below into {target_lang}, adhering to the following rules:
                    1. Keep the original tone and style, ensuring the translation is concise and suitable for presentation slides.
                    2. Avoid overly formal or verbose expressions. For example:
                       - In Japanese, avoid 'です' and 'ます' unless absolutely necessary.
                       - In Chinese, use straightforward and professional wording.
                       - In English, prioritize brevity and clarity.
                    3. Do not translate or modify placeholders in the format [PLACEHOLDER_X]. Ensure placeholders remain in their original positions and with the same quantity as in the original text.
                    4. Only output the translated text without explanations or extra comments.

                    Text: {text}
                """

                # Set other parameters
                chat_request.max_tokens = 2000
                chat_request.temperature = 0
                chat_request.frequency_penalty = 0
                chat_request.top_p = 0.75
                chat_request.top_k = 0
                chat_request.is_stream = False

                chat_detail.serving_mode = oci.generative_ai_inference.models.OnDemandServingMode(
                    model_id=model_name
                )
                chat_detail.chat_request = chat_request
                chat_detail.compartment_id = os.environ.get("COMPARTMENT_ID")

                # Make the API call
                chat_response = generative_ai_inference_client.chat(chat_detail)

                print(f"{text=}")
                print(f"translated: {chat_response.data.chat_response.text}\n")
                return chat_response.data.chat_response.text
            except Exception as e:
                if attempt < max_attempts - 1:
                    print(f"Attempt {attempt + 1} failed with error: {str(e)}. Retrying...")
                    time.sleep(1)
                else:
                    print(f"All {max_attempts} attempts failed with error: {str(e)}")
                    return text


def translate_ppt(model_name, input_ppt, target_lang):
    # 读取输入PPT
    ppt = Presentation(input_ppt)
    input_file_name = os.path.basename(input_ppt)

    # 遍历每一张幻灯片
    for slide_index, slide in enumerate(ppt.slides, start=1):
        print(f'Translate slide {slide_index}/{len(ppt.slides)}')
        print('-------------------------------------------')
        for shape in slide.shapes:
            # 跳过页脚部分
            if shape.is_placeholder and shape.placeholder_format.type in [
                PP_PLACEHOLDER_TYPE.FOOTER,
                PP_PLACEHOLDER_TYPE.SLIDE_NUMBER,
                PP_PLACEHOLDER_TYPE.DATE,
            ]:
                continue

            # print(f"{shape.shape_type=}")
            # print(f"{shape.has_table=}, {shape.has_text_frame=}")
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        original_text = cell.text_frame.text
                        if original_text and original_text.strip() and len(original_text.strip()) > 0:
                            translated_text = translate_text(original_text, target_lang, model_name)
                            cell.text_frame.text = translated_text
                            # 处理 SmartArt
            elif shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    # Step 1: 提取段落的完整文本，并为每个 run 添加唯一标记符
                    original_runs = []
                    full_text_with_delimiters = ""

                    for idx, run in enumerate(paragraph.runs):
                        original_text = run.text.strip()
                        if original_text and len(original_text) > 0:
                            delimiter = f"[PLACEHOLDER_{idx}]"  # 唯一标记符
                            full_text_with_delimiters += f"{delimiter}{original_text}"
                            original_runs.append({"run": run, "delimiter": delimiter})

                    if full_text_with_delimiters == "":
                        continue
                    # Step 2: 翻译整个段落（包含标记符）
                    translated_text_with_delimiters = translate_text(full_text_with_delimiters, target_lang, model_name)

                    # Step 3: 根据标记符分割翻译结果，并写回每个 run
                    for item in original_runs:
                        delimiter = item["delimiter"]
                        run = item["run"]

                        # 找到标记符的位置，并提取对应翻译文本
                        start_idx = translated_text_with_delimiters.find(delimiter)
                        if start_idx != -1:
                            end_idx = start_idx + len(delimiter)
                            # 提取翻译后的内容，去掉标记符
                            translated_run_text = translated_text_with_delimiters[end_idx:].split("[PLACEHOLDER_", 1)[0]
                            run.text = translated_run_text

                    # for run in paragraph.runs:
                    #     original_text = run.text
                    #     if original_text and original_text.strip():
                    #         translated_text = translate_text(original_text, target_lang, model_name)
                    #         run.text = translated_text

        if slide.has_notes_slide:
            original_text = slide.notes_slide.notes_text_frame.text
            if original_text and original_text.strip() and len(original_text.strip()) > 0:
                translated_text = translate_text(original_text, target_lang, model_name)
                slide.notes_slide.notes_text_frame.text = translated_text

    # 保存翻译后的PPT
    output_file_name = f"{input_file_name.rsplit('.', 1)[0]}_{target_lang}.{input_file_name.rsplit('.', 1)[-1]}"

    # Save the PowerPoint file
    output_file_path = os.path.join(output_dir, output_file_name)
    print(f"{output_file_path=}")
    ppt.save(output_file_path)
    print("Translation completed.")
    return output_file_path


model_type = gr.Radio(
    choices=[
        "cohere.command-r-08-2024",
        "cohere.command-r-plus-08-2024",
        "gpt-4"
    ],
    label="Model Type",
    value="gpt-4"
)
input_ppt = gr.File(label="Upload PPTX file", file_types=[".pptx"])
target_lang = gr.Dropdown(choices=["English", "Japanese", "Chinese"], label="Target Language", value="Japanese")
output_ppt = gr.File(label="Download Translated PPTX")

gr.Interface(
    fn=translate_ppt,
    inputs=[model_type, input_ppt, target_lang],
    outputs=output_ppt,
    title="PPT Translator",
    description="Upload a PPTX file and specify target language to get a translated PPTX file."
).launch(server_port=8080)
